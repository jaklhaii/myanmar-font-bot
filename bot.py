#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
⚡ Professional Myanmar Font Converter & Document Processing Bot
- Pyidaungsu Font Converter
- Advanced PDF Text Cleaning & Reordering (Visual to Logical)
- PowerPoint (.pptx) Slide Image Extractor with Interaction Flow
"""

import logging
import os
import sys
import tempfile
import unicodedata
import re
import shutil
import converter
import rabbit
from pypdf import PdfReader
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder,
    ContextTypes,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    filters,
)

# Configure logging format
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

# Retrieve Telegram Bot Token
TOKEN = os.getenv("TELEGRAM_BOT_TOKEN") or os.getenv("TELEGRAM_TOKEN") or "8183997269:AAHF5VgSgR7TJhC0HX9QgPCs74olBmoh2eA"
FONT_PATH = "Pyidaungsu.ttf"

def clean_myanmar_pdf_text(text: str) -> str:
    """
    Clean and convert Myanmar text. 
    Handles:
    1. Visual order to Logical order (common in PDF extraction)
    2. Zawgyi to Unicode conversion
    3. PUA character cleaning
    """
    if not text:
        return ""
    
    # Step 1: Remove PUA characters often used in legacy fonts
    text = re.sub(r'[\uE000-\uF8FF]', '', text)
    
    # Step 2: Normalize Unicode
    text = unicodedata.normalize('NFC', text)
    
    # Step 3: Fix Visual Order Unicode (Common PDF issue)
    # Rule: ေ (U+1031) and ြ (U+103C) often appear before the consonant in visual extraction
    # Reorder: [ေ/ြ] + Consonant -> Consonant + [ေ/ြ]
    # We use a loop to handle multiple combinations
    for _ in range(3):
        # ေ (1031) + Consonant (1000-1021)
        text = re.sub(r'(\u1031)([\u1000-\u1021])', r'\2\1', text)
        # ြ (103C) + Consonant (1000-1021)
        text = re.sub(r'(\u103c)([\u1000-\u1021])', r'\2\1', text)
        # ေ (1031) + Medials (103B-103E)
        text = re.sub(r'(\u1031)([\u103b-\u103e])', r'\2\1', text)

    # Step 4: Zawgyi to Unicode Conversion
    # We check for Zawgyi-specific patterns or visual order leftovers
    is_zawgyi = False
    # Zawgyi specific codepoints
    if re.search(r'[\u107e-\u1084\u1088\u1089\u1090\u1091\u1092\u1097]', text):
        is_zawgyi = True
    # If 'e' is still before consonant after our reordering attempts (shouldn't happen but for safety)
    elif re.search(r'\u1031[\u1000-\u1021]', text):
        is_zawgyi = True
        
    if is_zawgyi:
        try:
            text = rabbit.zg2uni(text)
        except Exception as e:
            logger.error(f"Rabbit conversion error: {e}")

    # Step 5: Universal Normalization via converter module
    try:
        text = converter.to_pyidaungsu(text)
    except Exception as e:
        logger.error(f"Converter module error: {e}")
        
    return text

async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Send a stunning pro-style welcome message."""
    user = update.effective_user
    welcome_text = (
        f"✨ မင်္ဂလာပါ {user.first_name} ခင်ဗျာ ✨\n\n"
        "──────────────────────────────\n"
        "🎯 Pyidaungsu Font & Document Bot Pro\n"
        "──────────────────────────────\n\n"
        "📌 လုပ်ဆောင်ချက်များ:\n"
        "၁။ Font Converter: မည်သည့် မြန်မာစာသားကိုမဆို Pyidaungsu Unicode သို့ ပြောင်းပေးခြင်း။\n"
        "၂။ PDF Text Extractor: PDF ဖိုင်ပို့ပါက စာများဖတ်၍ ပြန်ပို့ပေးခြင်း။\n"
        "၃။ PPTX to Images: PowerPoint ဖိုင်ပို့ပါက Slide များကို ပုံများအဖြစ် ပြောင်းပေးခြင်း။\n\n"
        "💡 အသုံးပြုရန် ဖိုင် သို့မဟုတ် စာသားများကို ပို့ပေးပါ။"
    )

    keyboard = [
        [InlineKeyboardButton("📖 အသုံးပြုပုံ လမ်းညွှန်", callback_data="help")],
        [InlineKeyboardButton("💎 Pro Features", callback_data="features")]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)

    if update.message:
        await update.message.reply_text(welcome_text, reply_markup=reply_markup)
    elif update.callback_query:
        await update.callback_query.message.edit_text(welcome_text, reply_markup=reply_markup)

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Provide guide."""
    help_text = (
        "📖 အသုံးပြုပုံ လမ်းညွှန်\n\n"
        "• စာသားများ: ဇော်ဂျီ သို့မဟုတ် ဖောင့်မမှန်သည်များကို ပို့ပါက Pyidaungsu သို့ ပြောင်းပေးမည်။\n"
        "• PDF/PPTX ဖိုင်များ: ဖိုင်ပို့ပါက ပထမစာမျက်နှာကို အရင်ပြသမည်ဖြစ်ပြီး ကျန်ရှိသည်များကို တစ်ခုချင်းစီ သို့မဟုတ် အားလုံးကို တစ်ခါတည်း ထုတ်ယူနိုင်ပါသည်။"
    )
    query = update.callback_query
    if query:
        await query.answer()
        keyboard = [[InlineKeyboardButton("⬅️ ပင်မမီနူးသို့", callback_data="main_menu")]]
        await query.message.edit_text(help_text, reply_markup=InlineKeyboardMarkup(keyboard))
    else:
        await update.message.reply_text(help_text)

async def features_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Highlight features."""
    features_text = (
        "💎 Pro Features များနှင့် အားသာချက်များ\n\n"
        "• Smart Font Detection & Conversion\n"
        "• Interactive Page-by-Page Extraction\n"
        "• High Quality Myanmar Text Rendering (Pyidaungsu)"
    )
    query = update.callback_query
    if query:
        await query.answer()
        keyboard = [[InlineKeyboardButton("⬅️ ပင်မမီနူးသို့", callback_data="main_menu")]]
        await query.message.edit_text(features_text, reply_markup=InlineKeyboardMarkup(keyboard))

def render_pptx_slide(input_path, slide_idx, total_slides):
    """Render a single PPTX slide to an image."""
    prs = Presentation(input_path)
    if slide_idx >= len(prs.slides):
        return None
    
    slide = prs.slides[slide_idx]
    img = Image.new('RGB', (1280, 720), color=(245, 247, 250))
    draw = ImageDraw.Draw(img)
    
    # Load Font
    try:
        font_title = ImageFont.truetype(FONT_PATH, 42)
        font_body = ImageFont.truetype(FONT_PATH, 30)
    except:
        font_title = font_body = None

    slide_text_lines = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                p_text = paragraph.text.strip()
                if p_text:
                    slide_text_lines.append(clean_myanmar_pdf_text(p_text))

    draw.rectangle([0, 0, 1280, 100], fill=(24, 43, 73))
    draw.text((50, 30), f"Slide {slide_idx + 1} / {total_slides}", fill=(255, 255, 255), font=font_title)
    
    y_offset = 140
    for line in slide_text_lines[:15]:
        draw.text((60, y_offset), line[:70], fill=(30, 30, 30), font=font_body)
        y_offset += 45

    temp_img_path = os.path.join(tempfile.gettempdir(), f"slide_render_{slide_idx}.png")
    img.save(temp_img_path, 'PNG')
    return temp_img_path

def extract_pdf_page(input_path, page_idx):
    """Extract and clean text from a single PDF page."""
    reader = PdfReader(input_path)
    if page_idx >= len(reader.pages):
        return None
    
    text = reader.pages[page_idx].extract_text()
    if text:
        cleaned = clean_myanmar_pdf_text(text)
        return f"--- Page {page_idx + 1} ---\n{cleaned}"
    return f"--- Page {page_idx + 1} ---\n(စာသားမတွေ့ရှိပါ)"

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle initial document upload."""
    message = update.message
    if not message or not message.document:
        return

    doc = message.document
    file_name = doc.file_name.lower()
    
    if not (file_name.endswith('.pdf') or file_name.endswith('.pptx')):
        await message.reply_text("❌ ကျေးဇူးပြု၍ **PDF** သို့မဟုတ် **PowerPoint (.pptx)** ဖိုင်များကိုသာ ပို့ပေးပါ။")
        return

    status_msg = await message.reply_text("⏳ ဖိုင်ကို လက်ခံရရှိပါပြီ။ စစ်ဆေးနေပါသည်...")

    try:
        file_obj = await doc.get_file()
        user_dir = os.path.join(tempfile.gettempdir(), f"user_{update.effective_user.id}")
        os.makedirs(user_dir, exist_ok=True)
        input_path = os.path.join(user_dir, doc.file_name)
        await file_obj.download_to_drive(input_path)

        doc_type = 'pdf' if file_name.endswith('.pdf') else 'pptx'
        total_pages = 0
        
        if doc_type == 'pdf':
            reader = PdfReader(input_path)
            total_pages = len(reader.pages)
        else:
            prs = Presentation(input_path)
            total_pages = len(prs.slides)

        context.user_data['current_doc'] = {
            'path': input_path,
            'type': doc_type,
            'total': total_pages,
            'current_idx': 0
        }

        await status_msg.delete()
        await send_next_part(update, context)

    except Exception as e:
        logger.error(f"Upload error: {e}")
        await message.reply_text(f"❌ အမှားအယွင်း ဖြစ်သွားပါသည်: {str(e)}")

async def send_next_part(update: Update, context: ContextTypes.DEFAULT_TYPE, send_all=False) -> None:
    """Process and send the next page or all pages."""
    doc_info = context.user_data.get('current_doc')
    if not doc_info:
        return

    idx = doc_info['current_idx']
    total = doc_info['total']
    path = doc_info['path']
    dtype = doc_info['type']

    while idx < total:
        if dtype == 'pdf':
            content = extract_pdf_page(path, idx)
            quoted = f"> {content.replace(chr(10), chr(10) + '> ')}"
            if update.callback_query:
                await update.callback_query.message.reply_text(quoted)
            else:
                await update.message.reply_text(quoted)
        else:
            img_path = render_pptx_slide(path, idx, total)
            caption = f"> 📄 Slide {idx + 1} / {total}"
            with open(img_path, 'rb') as photo:
                if update.callback_query:
                    await update.callback_query.message.reply_photo(photo=photo, caption=caption)
                else:
                    await update.message.reply_photo(photo=photo, caption=caption)

        idx += 1
        doc_info['current_idx'] = idx
        
        if not send_all:
            break

    if idx < total:
        keyboard = [
            [
                InlineKeyboardButton("➡️ နောက်တစ်မျက်နှာ", callback_data="next_page"),
                InlineKeyboardButton("⏩ အကုန်ထုတ်မယ်", callback_data="extract_all")
            ]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        text = f"စာမျက်နှာ {idx} ကို ပို့ဆောင်ပြီးပါပြီ။ ဆက်လက်ထုတ်ယူမလား?"
        
        if update.callback_query:
            await update.callback_query.message.reply_text(text, reply_markup=reply_markup)
        else:
            await update.message.reply_text(text, reply_markup=reply_markup)
    else:
        text = "✅ ဖိုင်တစ်ခုလုံး လုပ်ဆောင်ပြီးစီးပါပြီ။"
        if update.callback_query:
            await update.callback_query.message.reply_text(text)
        else:
            await update.message.reply_text(text)
        # Cleanup
        if os.path.exists(path):
            try: shutil.rmtree(os.path.dirname(path))
            except: pass
        context.user_data['current_doc'] = None

async def button_callback_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle interactive buttons."""
    query = update.callback_query
    await query.answer()

    if query.data == "help":
        await help_command(update, context)
    elif query.data == "features":
        await features_command(update, context)
    elif query.data == "main_menu":
        await start_command(update, context)
    elif query.data == "next_page":
        await send_next_part(update, context, send_all=False)
    elif query.data == "extract_all":
        await send_next_part(update, context, send_all=True)

async def handle_incoming_message(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle plain text conversion."""
    if not update.message or not update.message.text:
        return
    raw_text = update.message.text
    try:
        converted = clean_myanmar_pdf_text(raw_text)
        quoted = f"> {converted.replace(chr(10), chr(10) + '> ')}"
        await update.message.reply_text(quoted)
    except Exception as e:
        logger.error(f"Text error: {e}")
        await update.message.reply_text("❌ စာသားပြောင်းလဲရာတွင် အမှားရှိနေပါသည်။")

def main() -> None:
    """Initialize and run the Bot."""
    if not TOKEN or TOKEN == "YOUR_BOT_TOKEN_HERE":
        logger.critical("Bot Token is not configured.")
        sys.exit(1)

    application = ApplicationBuilder().token(TOKEN).build()

    application.add_handler(CommandHandler("start", start_command))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CallbackQueryHandler(button_callback_handler))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    application.add_handler(MessageHandler(filters.TEXT & (~filters.COMMAND), handle_incoming_message))

    logger.info("Enhanced Document Processing Bot is running...")
    application.run_polling()

if __name__ == "__main__":
    main()
