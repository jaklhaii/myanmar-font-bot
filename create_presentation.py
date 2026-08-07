import sys
import os
import converter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def main():
    prs = Presentation()
    
    # Add a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Raw text to be processed by converter (simulating bot logic)
    raw_title = "ပညာရေးနှင့် ခေတ်မီဖွံ့ဖြိုးတိုးတက်သော လူငယ်များ"
    raw_content = (
        "၁။ နိဒါန်း (Introduction)\n"
        "• ခေတ်မီဖွံ့ဖြိုးတိုးတက်သော နိုင်ငံတော်သစ်ကြီး တည်ဆောက်ရာတွင် လူငယ်တို့၏ အခန်းကဏ္ဍသည် အလွန်အရေးကြီးပါသည်။\n"
        "• ပညာရေးသည် လူငယ်များ၏ အနာဂတ်ကို ဖန်တီးပေးသည့် အဓိကသော့ချက် ဖြစ်ပါသည်။\n\n"
        "၂။ ပညာရေး၏ အခန်းကဏ္ဍ (Role of Education)\n"
        "• စာပေပညာသာမက နည်းပညာနှင့် ဗဟုသုတများကိုပါ စဉ်ဆက်မပြတ် ဆည်းပူးသင်ယူရမည် ဖြစ်ပါသည်။\n"
        "• အသိပညာ၊ အတတ်ပညာနှင့် ပြည့်စုံသော လူငယ်များပေါ်ထွက်လာမှသာ နိုင်ငံတော် တိုးတက်မည် ဖြစ်ပါသည်။\n\n"
        "၃။ နိဂုံး (Conclusion)\n"
        "• လူငယ်တို့သည် နိုင်ငံတော်၏ အနာဂတ် ရတနာများ ဖြစ်ကြသည်နှင့်အညီ ဗလငါးတန်နှင့် ပြည့်စုံသော သားကောင်းရတနာများ ဖြစ်အောင် ကြိုးစားကြပါစို့။"
    )
    
    # Process text using bot converter logic
    title_text = converter.to_pyidaungsu(raw_title)
    content_text = converter.to_pyidaungsu(raw_content)
    
    # Add Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.5), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Pyidaungsu"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(24, 43, 73)
    p_title.alignment = PP_ALIGN.CENTER
    
    # Add Content Box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(4.8))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    paragraphs = content_text.split('\n')
    for i, line in enumerate(paragraphs):
        if i == 0:
            p = tf_content.paragraphs[0]
        else:
            p = tf_content.add_paragraph()
        p.text = line
        p.font.name = "Pyidaungsu"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(6)

    output_path = "/home/ubuntu/bot_extracted/Essay_Pyidaungsu.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    main()
